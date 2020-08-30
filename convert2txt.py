import docx2txt
import pptx

import glob
import os

from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage

def extract_txt_from_docx(file_name):
    text = docx2txt.process(file_name)

    line_list = list()
    for line in text:
        line_list.append(line)
        #f.write(line)

    prepro_line = list()
    del_count = 0
    for i, line in enumerate(line_list):
        if i < 3:
            prepro_line.append(line)
            continue
        else:
            if prepro_line[i-1-del_count] == "\n" and prepro_line[i-2-del_count] == "\n":
                del prepro_line[-1]
                del_count += 1
                prepro_line.append(line)
            else:
                prepro_line.append(line)

    return prepro_line

def write_docxfile(file_list):
    for i, doc_file in enumerate(file_list):
        file_txt = extract_txt_from_docx(doc_file)
        #print(file_txt)
        if i == 0:
            with open("./all_docx.txt", "w") as f:
                for line in file_txt:
                    f.write(line)
        else:
            with open("./all_docx.txt", "a") as f:
                f.write("\n\n")
                for line in file_txt:
                    f.write(line)
        print(doc_file, "had exported")


def extract_txt_from_pptx(file_name):
    #f = "./file_list/sample_pptx.pptx"
    for file in file_name:
        prs = pptx.Presentation(file_name)
        line_list = list()

        for i, sld in enumerate(prs.slides, start=1):
            #print(i, sld)

            for shp in sld.shapes:
                if shp.has_text_frame:
                    #print(shp.text)
                    line_list.append(shp.text)
        #print(line_list)
        return line_list

def write_pptxfile(file_list):
    for i, pptx_file in enumerate(file_list):
        file_txt = extract_txt_from_pptx(pptx_file)
        #print(file_txt)
        if i == 0:
            with open("./all_pptx.txt", "w") as f:
                for line in file_txt:
                    if line == "": continue
                    f.writelines(line.replace("\n",""))
                    f.write("\n")
        else:
            with open("./all_pptx.txt", "a") as f:
                f.write("\n\n")
                for line in file_txt:
                    if line == "": continue
                    f.writelines(line.replace("\n",""))
                    f.write("\n")
        print(pptx_file, "had exported")


def extract_txt_from_pdf(file_name):
    input_path = file_name
    output_path = 'result.txt' # 一時ファイル

    manager = PDFResourceManager()

    with open(output_path, "wb") as output:
        with open(input_path, 'rb') as input:
            with TextConverter(manager, output, codec='utf-8', laparams=LAParams()) as conv:
                interpreter = PDFPageInterpreter(manager, conv)
                for page in PDFPage.get_pages(input):
                    interpreter.process_page(page)

    with open("result.txt", encoding="utf-8") as f:
        line_list = list()
        for line in f:
            if line == "" or line == "\n" or line == " \n":
                continue
            else:
                line_list.append(line)
    
    os.remove("./result.txt") #一時ファイルの削除
    #print(line_list)

    return line_list

def write_pdffile(file_list):
    for i, pdf_file in enumerate(file_list):
        file_txt = extract_txt_from_pdf(pdf_file)
        #print(file_txt)
        if i == 0:
            with open("./all_pdf.txt", "w") as f:
                for line in file_txt:
                    if line == "": continue
                    f.writelines(line.replace("\n",""))
                    f.write("\n")
        else:
            with open("./all_pdf.txt", "a") as f:
                f.write("\n\n")
                for line in file_txt:
                    if line == "": continue
                    f.writelines(line.replace("\n",""))
                    f.write("\n")
        print(pdf_file, "had exported")


if __name__ == "__main__":
    doc_file_list = glob.glob("./file_list/*.docx")
    if len(doc_file_list) > 0:
        write_docxfile(doc_file_list)
    else:
        print(".docx file is none :/")

    pptx_file_list = glob.glob("./file_list/*.pptx")
    if len(pptx_file_list) > 0:
        write_pptxfile(pptx_file_list)
    else:
        print(".pptx file is none :/")

    pdf_file_list = glob.glob("./file_list/*.pdf")
    if len(pdf_file_list) > 0:
        write_pdffile(pdf_file_list)
    else:
        print(".pdf file is none :/")
